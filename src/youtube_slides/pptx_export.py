from __future__ import annotations

import html
import shutil
import tempfile
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from .pipeline import SlideInfo

PML = "http://schemas.openxmlformats.org/presentationml/2006/main"
DML = "http://schemas.openxmlformats.org/drawingml/2006/main"
REL = "http://schemas.openxmlformats.org/package/2006/relationships"
OD_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
CT = "http://schemas.openxmlformats.org/package/2006/content-types"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

ET.register_namespace("p", PML)
ET.register_namespace("a", DML)
ET.register_namespace("r", R_NS)
ET.register_namespace("", CT)


def export_pptx(
    slides_dir: str | Path,
    slides: list[SlideInfo],
    output_path: str | Path,
    *,
    notes: dict[int, list[str] | str] | None = None,
) -> Path:
    """Create a 16:9 PowerPoint deck from extracted slide PNGs."""

    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception as exc:
        raise RuntimeError("python-pptx is required to export PowerPoint files") from exc

    slides_dir = Path(slides_dir)
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    blank_layout = presentation.slide_layouts[6]

    for slide_info in slides:
        slide = presentation.slides.add_slide(blank_layout)
        image_path = slides_dir / slide_info.file
        slide.shapes.add_picture(
            str(image_path),
            0,
            0,
            width=presentation.slide_width,
            height=presentation.slide_height,
        )

    presentation.save(output_path)

    if notes:
        add_speaker_notes(output_path, slides, notes)
    return output_path


def write_reading_view(
    slides_dir: str | Path,
    slides: list[SlideInfo],
    output_path: str | Path,
    *,
    notes: dict[int, list[str] | str] | None = None,
) -> Path:
    slides_dir = Path(slides_dir)
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    figures = []
    for slide in slides:
        note_value = notes.get(slide.index, "") if notes else ""
        note_lines = note_value if isinstance(note_value, list) else [note_value]
        note_html = "".join(f"<p>{html.escape(line)}</p>" for line in note_lines if line)
        figures.append(
            "<section>"
            f'<img src="{html.escape(str((slides_dir / slide.file).resolve()))}" alt="Slide {slide.index}">'
            f"<h2>Slide {slide.index:03d} <span>{html.escape(slide.timestamp_label)}</span></h2>"
            f"{note_html}"
            "</section>"
        )

    document = f"""<!doctype html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>YouTube lecture slides reading view</title>
<style>
body {{ margin: 0; font: 16px/1.55 -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif; color: #172033; background: #f5f6f8; }}
main {{ max-width: 1040px; margin: 0 auto; padding: 32px 18px 56px; }}
section {{ margin: 0 0 28px; background: white; border: 1px solid #d8dee6; border-radius: 8px; overflow: hidden; }}
img {{ width: 100%; display: block; background: #eef1f5; }}
h2 {{ margin: 18px 22px 8px; font-size: 20px; }}
h2 span {{ color: #687385; font-weight: 500; margin-left: 8px; }}
p {{ margin: 8px 22px 18px; }}
</style>
</head>
<body>
<main>
{''.join(figures)}
</main>
</body>
</html>
"""
    output_path.write_text(document, encoding="utf-8")
    return output_path


def add_speaker_notes(
    pptx_path: str | Path,
    slides: list[SlideInfo],
    notes: dict[int, list[str] | str],
) -> None:
    pptx_path = Path(pptx_path)
    with tempfile.TemporaryDirectory(prefix="youtube-slides-pptx-") as tmp:
        tmp_path = Path(tmp)
        source = tmp_path / "source.pptx"
        target = tmp_path / "target.pptx"
        shutil.copy2(pptx_path, source)
        with zipfile.ZipFile(source, "r") as zin, zipfile.ZipFile(target, "w", zipfile.ZIP_DEFLATED) as zout:
            names = set(zin.namelist())
            for item in zin.infolist():
                if _is_generated_notes_part(item.filename):
                    continue
                data = zin.read(item.filename)
                if item.filename == "[Content_Types].xml":
                    data = _update_content_types(data, len(slides))
                elif item.filename == "ppt/presentation.xml":
                    data = _update_presentation_xml(data)
                elif item.filename == "ppt/_rels/presentation.xml.rels":
                    data = _update_presentation_rels(data)
                elif item.filename.startswith("ppt/slides/_rels/slide") and item.filename.endswith(".xml.rels"):
                    slide_number = _slide_number_from_rels_path(item.filename)
                    data = _update_slide_rels(data, slide_number)
                zout.writestr(item, data)

            zout.writestr("ppt/notesMasters/notesMaster1.xml", _notes_master_xml())
            zout.writestr("ppt/notesMasters/_rels/notesMaster1.xml.rels", _notes_master_rels_xml())
            zout.writestr("ppt/notesMasters/theme/themeNotes.xml", _notes_theme_xml())
            for offset, slide in enumerate(slides, start=1):
                note_value = notes.get(slide.index, "")
                note_lines = note_value if isinstance(note_value, list) else [note_value]
                zout.writestr(f"ppt/notesSlides/notesSlide{offset}.xml", _notes_slide_xml(note_lines))
                zout.writestr(
                    f"ppt/notesSlides/_rels/notesSlide{offset}.xml.rels",
                    _notes_slide_rels_xml(offset),
                )
        shutil.move(target, pptx_path)


def _is_generated_notes_part(name: str) -> bool:
    return (
        name.startswith("ppt/notesSlides/")
        or name.startswith("ppt/notesMasters/")
    )


def _update_content_types(data: bytes, slide_count: int) -> bytes:
    root = ET.fromstring(data)
    existing = {child.attrib.get("PartName") for child in root}
    overrides = [
        ("/ppt/notesMasters/notesMaster1.xml", "application/vnd.openxmlformats-officedocument.presentationml.notesMaster+xml"),
        ("/ppt/notesMasters/theme/themeNotes.xml", "application/vnd.openxmlformats-officedocument.theme+xml"),
    ]
    for index in range(1, slide_count + 1):
        overrides.append(
            (
                f"/ppt/notesSlides/notesSlide{index}.xml",
                "application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml",
            )
        )
    for part_name, content_type in overrides:
        if part_name not in existing:
            ET.SubElement(root, f"{{{CT}}}Override", PartName=part_name, ContentType=content_type)
    return _xml_bytes(root)


def _update_presentation_xml(data: bytes) -> bytes:
    root = ET.fromstring(data)
    if root.find(f"{{{PML}}}notesMasterIdLst") is not None:
        return data
    notes_list = ET.Element(f"{{{PML}}}notesMasterIdLst")
    ET.SubElement(notes_list, f"{{{PML}}}notesMasterId", {f"{{{R_NS}}}id": "rIdNotesMaster"})
    insert_at = 1 if len(root) > 0 and root[0].tag == f"{{{PML}}}sldMasterIdLst" else 0
    root.insert(insert_at, notes_list)
    return _xml_bytes(root)


def _update_presentation_rels(data: bytes) -> bytes:
    root = ET.fromstring(data)
    for relationship in root:
        if relationship.attrib.get("Type") == f"{OD_REL}/notesMaster":
            return data
    ET.SubElement(
        root,
        f"{{{REL}}}Relationship",
        Id="rIdNotesMaster",
        Type=f"{OD_REL}/notesMaster",
        Target="notesMasters/notesMaster1.xml",
    )
    return _xml_bytes(root)


def _update_slide_rels(data: bytes, slide_number: int) -> bytes:
    root = ET.fromstring(data)
    for relationship in root:
        if relationship.attrib.get("Type") == f"{OD_REL}/notesSlide":
            return data
    ET.SubElement(
        root,
        f"{{{REL}}}Relationship",
        Id="rIdNotesSlide",
        Type=f"{OD_REL}/notesSlide",
        Target=f"../notesSlides/notesSlide{slide_number}.xml",
    )
    return _xml_bytes(root)


def _slide_number_from_rels_path(path: str) -> int:
    stem = Path(path).name
    return int(stem.removeprefix("slide").removesuffix(".xml.rels"))


def _notes_slide_xml(lines: list[str]) -> str:
    paragraphs = []
    for line in lines:
        if not line:
            continue
        paragraphs.append(
            f'<a:p><a:r><a:rPr lang="en-US"/><a:t>{html.escape(line)}</a:t></a:r>'
            '<a:endParaRPr lang="en-US"/></a:p>'
        )
    if not paragraphs:
        paragraphs.append('<a:p><a:endParaRPr lang="en-US"/></a:p>')
    body = "".join(paragraphs)
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:notes xmlns:a="{DML}" xmlns:r="{R_NS}" xmlns:p="{PML}">
  <p:cSld>
    <p:spTree>
      <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
      <p:grpSpPr><a:xfrm/></p:grpSpPr>
      <p:sp>
        <p:nvSpPr><p:cNvPr id="2" name="Slide Image Placeholder 1"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph type="sldImg" idx="0"/></p:nvPr></p:nvSpPr>
        <p:spPr/>
      </p:sp>
      <p:sp>
        <p:nvSpPr><p:cNvPr id="3" name="Notes Placeholder 2"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph type="body" idx="1"/></p:nvPr></p:nvSpPr>
        <p:spPr/>
        <p:txBody><a:bodyPr/><a:lstStyle/>{body}</p:txBody>
      </p:sp>
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:notes>'''


def _notes_slide_rels_xml(slide_number: int) -> str:
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="{REL}">
  <Relationship Id="rIdSlide" Type="{OD_REL}/slide" Target="../slides/slide{slide_number}.xml"/>
  <Relationship Id="rIdNotesMaster" Type="{OD_REL}/notesMaster" Target="../notesMasters/notesMaster1.xml"/>
</Relationships>'''


def _notes_master_xml() -> str:
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:notesMaster xmlns:a="{DML}" xmlns:r="{R_NS}" xmlns:p="{PML}">
  <p:cSld>
    <p:bg><p:bgRef idx="1001"><a:schemeClr val="bg1"/></p:bgRef></p:bg>
    <p:spTree>
      <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
      <p:grpSpPr><a:xfrm/></p:grpSpPr>
      <p:sp><p:nvSpPr><p:cNvPr id="2" name="Header Placeholder"/><p:cNvSpPr/><p:nvPr><p:ph type="hdr" sz="quarter"/></p:nvPr></p:nvSpPr><p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>
      <p:sp><p:nvSpPr><p:cNvPr id="3" name="Date Placeholder"/><p:cNvSpPr/><p:nvPr><p:ph type="dt" sz="quarter" idx="1"/></p:nvPr></p:nvSpPr><p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>
      <p:sp><p:nvSpPr><p:cNvPr id="4" name="Slide Image Placeholder"/><p:cNvSpPr/><p:nvPr><p:ph type="sldImg" idx="2"/></p:nvPr></p:nvSpPr><p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>
      <p:sp><p:nvSpPr><p:cNvPr id="5" name="Notes Placeholder"/><p:cNvSpPr/><p:nvPr><p:ph type="body" idx="3"/></p:nvPr></p:nvSpPr><p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:defRPr sz="1200"/></a:p></p:txBody></p:sp>
      <p:sp><p:nvSpPr><p:cNvPr id="6" name="Footer Placeholder"/><p:cNvSpPr/><p:nvPr><p:ph type="ftr" sz="quarter" idx="4"/></p:nvPr></p:nvSpPr><p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>
      <p:sp><p:nvSpPr><p:cNvPr id="7" name="Slide Number Placeholder"/><p:cNvSpPr/><p:nvPr><p:ph type="sldNum" sz="quarter" idx="5"/></p:nvPr></p:nvSpPr><p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>
    </p:spTree>
  </p:cSld>
  <p:clrMap bg1="lt1" tx1="dk1" bg2="lt2" tx2="dk2" accent1="accent1" accent2="accent2" accent3="accent3" accent4="accent4" accent5="accent5" accent6="accent6" hlink="hlink" folHlink="folHlink"/>
</p:notesMaster>'''


def _notes_master_rels_xml() -> str:
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="{REL}">
  <Relationship Id="rIdTheme" Type="{OD_REL}/theme" Target="theme/themeNotes.xml"/>
</Relationships>'''


def _notes_theme_xml() -> str:
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<a:theme xmlns:a="{DML}" name="Office">
  <a:themeElements>
    <a:clrScheme name="Office">
      <a:dk1><a:srgbClr val="000000"/></a:dk1><a:lt1><a:srgbClr val="FFFFFF"/></a:lt1>
      <a:dk2><a:srgbClr val="1F497D"/></a:dk2><a:lt2><a:srgbClr val="EEECE1"/></a:lt2>
      <a:accent1><a:srgbClr val="4F81BD"/></a:accent1><a:accent2><a:srgbClr val="C0504D"/></a:accent2>
      <a:accent3><a:srgbClr val="9BBB59"/></a:accent3><a:accent4><a:srgbClr val="8064A2"/></a:accent4>
      <a:accent5><a:srgbClr val="4BACC6"/></a:accent5><a:accent6><a:srgbClr val="F79646"/></a:accent6>
      <a:hlink><a:srgbClr val="0000FF"/></a:hlink><a:folHlink><a:srgbClr val="800080"/></a:folHlink>
    </a:clrScheme>
    <a:fontScheme name="Office"><a:majorFont><a:latin typeface="Calibri Light"/></a:majorFont><a:minorFont><a:latin typeface="Calibri"/></a:minorFont></a:fontScheme>
    <a:fmtScheme name="Office"><a:fillStyleLst><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:fillStyleLst><a:lnStyleLst><a:ln w="6350"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln></a:lnStyleLst><a:effectStyleLst><a:effectStyle><a:effectLst/></a:effectStyle></a:effectStyleLst><a:bgFillStyleLst><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:bgFillStyleLst></a:fmtScheme>
  </a:themeElements>
</a:theme>'''


def _xml_bytes(root: ET.Element) -> bytes:
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)
